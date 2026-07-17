#!/usr/bin/env python3
"""
ERCOT text extraction pipeline — step 5.

Reads documents from ercot.db, extracts text from each file in docs/,
stores results in document_text table, then builds the FTS5 index.

Supported formats:
  pptx  → python-pptx (slides + speaker notes)
  pdf   → PyMuPDF
  docx  → python-docx
  zip   → tries pptx, then docx (most zips are renamed Office files)
  xlsx  → openpyxl
  xls   → xlrd
  html  → beautifulsoup4
  txt   → plain read
  doc/ppt → skipped (need LibreOffice; marked 'skipped:legacy')

Resumable: skips documents where extraction_status is already set.

Usage:
  pip install python-pptx PyMuPDF python-docx openpyxl xlrd beautifulsoup4
  python extract_text.py
  python extract_text.py --limit 20     # test run
  python extract_text.py --redo-failed  # retry previously failed docs
"""

import argparse
import sqlite3
import traceback
from pathlib import Path

DOCS_DIR = Path("docs")
DB_PATH  = Path("ercot.db")


# ── Extractors ────────────────────────────────────────────────────────────────

def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(notes)
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    import fitz
    doc = fitz.open(str(path))
    parts = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            parts.append(text.strip())
    doc.close()
    return "\n".join(parts)


def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            row_text = " ".join(str(v) for v in row if v is not None).strip()
            if row_text:
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)


def extract_xls(path: Path) -> str:
    import xlrd
    wb = xlrd.open_workbook(str(path))
    parts = []
    for sheet in wb.sheets():
        for row_idx in range(sheet.nrows):
            row_text = " ".join(str(v) for v in sheet.row_values(row_idx) if v).strip()
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def extract_zip(path: Path) -> str:
    import zipfile

    # First try as Office formats (docx/pptx/xlsx are ZIP-based)
    for fn in (extract_pptx, extract_docx, extract_xlsx):
        try:
            text = fn(path)
            if text.strip():
                return text
        except Exception:
            continue

    # Treat as a real ZIP archive — extract text from each contained file
    parts = []
    try:
        with zipfile.ZipFile(path) as zf:
            for name in zf.namelist():
                if name.endswith("/"):
                    continue  # skip directories
                ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
                extractor = EXTRACTORS.get(ext)
                if not extractor or ext == "zip":
                    continue
                tmp = path.parent / f"_tmp_{path.stem}_{name.replace('/', '_')}"
                try:
                    tmp.write_bytes(zf.read(name))
                    text = extractor(tmp).strip()
                    if text:
                        parts.append(f"[{name}]\n{text}")
                except Exception:
                    continue
                finally:
                    if tmp.exists():
                        tmp.unlink()
    except Exception:
        pass

    return "\n\n".join(parts)


def extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(path.read_bytes(), "html.parser")
    return soup.get_text(" ", strip=True)


def extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


EXTRACTORS = {
    "pptx": extract_pptx,
    "pdf":  extract_pdf,
    "docx": extract_docx,
    "xlsx": extract_xlsx,
    "xls":  extract_xls,
    "zip":  extract_zip,
    "html": extract_html,
    "txt":  extract_txt,
}

SKIP_TYPES = {"doc", "ppt", "ole"}  # need LibreOffice; deferred


# ── Schema migration ──────────────────────────────────────────────────────────

SCHEMA = """
CREATE TABLE IF NOT EXISTS document_text (
    document_id INTEGER PRIMARY KEY REFERENCES documents(id),
    text        TEXT NOT NULL DEFAULT ''
);

DROP TABLE IF EXISTS document_fts;
CREATE VIRTUAL TABLE document_fts USING fts5(
    text,
    content='document_text',
    content_rowid='document_id'
);
"""


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description="Extract text from ERCOT documents into ercot.db.")
    ap.add_argument("--limit",       type=int, default=0, help="Process only first N docs.")
    ap.add_argument("--redo-failed", action="store_true", help="Retry docs marked failed.")
    ap.add_argument("--redo-empty",  action="store_true", help="Retry docs that extracted as empty.")
    ap.add_argument("--reindex",     action="store_true", help="Re-extract docs marked done but missing from document_text.")
    args = ap.parse_args()

    db = sqlite3.connect(DB_PATH)
    db.row_factory = sqlite3.Row

    print("Setting up document_text + document_fts tables...")
    db.executescript(SCHEMA)
    db.commit()

    # Fetch docs to process
    if args.redo_failed:
        rows = db.execute(
            "SELECT * FROM documents WHERE extraction_status LIKE 'failed:%'"
        ).fetchall()
    elif args.redo_empty:
        rows = db.execute(
            "SELECT * FROM documents WHERE extraction_status = 'empty'"
        ).fetchall()
    elif args.reindex:
        rows = db.execute(
            """SELECT * FROM documents
               WHERE extraction_status = 'done'
               AND id NOT IN (SELECT document_id FROM document_text)"""
        ).fetchall()
    else:
        rows = db.execute(
            "SELECT * FROM documents WHERE extraction_status = ''"
        ).fetchall()

    if args.limit:
        rows = rows[: args.limit]

    total = len(rows)
    print(f"{total} documents to process\n")

    n_ok = n_skipped = n_failed = n_empty = 0
    type_stats: dict[str, int] = {}

    for i, row in enumerate(rows, 1):
        ext      = row["url_ext"].lower()
        filename = row["filename"]
        doc_id   = row["id"]
        path     = DOCS_DIR / filename

        if ext in SKIP_TYPES:
            db.execute(
                "UPDATE documents SET extraction_status=? WHERE id=?",
                ("skipped:legacy", doc_id),
            )
            n_skipped += 1
            if i % 50 == 0 or i == total:
                print(f"  {i}/{total} | ok={n_ok} skip={n_skipped} fail={n_failed} empty={n_empty}")
            continue

        if not path.exists():
            db.execute(
                "UPDATE documents SET extraction_status=? WHERE id=?",
                ("failed:file_missing", doc_id),
            )
            n_failed += 1
            continue

        extractor = EXTRACTORS.get(ext)
        if extractor is None:
            db.execute(
                "UPDATE documents SET extraction_status=? WHERE id=?",
                (f"skipped:unknown_type_{ext}", doc_id),
            )
            n_skipped += 1
            continue

        try:
            text = extractor(path)
        except Exception as exc:
            short = str(exc)[:120].replace("\n", " ")
            db.execute(
                "UPDATE documents SET extraction_status=? WHERE id=?",
                (f"failed:{short}", doc_id),
            )
            n_failed += 1
            if i % 50 == 0 or i == total:
                print(f"  {i}/{total} | ok={n_ok} skip={n_skipped} fail={n_failed} empty={n_empty}")
            continue

        if not text.strip():
            n_empty += 1
            db.execute(
                "UPDATE documents SET extraction_status=? WHERE id=?",
                ("empty", doc_id),
            )
        else:
            db.execute(
                "INSERT OR REPLACE INTO document_text (document_id, text) VALUES (?, ?)",
                (doc_id, text),
            )
            db.execute(
                "UPDATE documents SET extraction_status=? WHERE id=?",
                ("done", doc_id),
            )
            n_ok += 1
            type_stats[ext] = type_stats.get(ext, 0) + 1

        if i % 50 == 0 or i == total:
            db.commit()
            print(f"  {i}/{total} | ok={n_ok} skip={n_skipped} fail={n_failed} empty={n_empty}")

    db.commit()

    print("\nRebuilding FTS index...")
    db.execute("INSERT INTO document_fts(document_fts) VALUES('rebuild')")
    db.commit()

    n_indexed = db.execute("SELECT COUNT(*) FROM document_text").fetchone()[0]
    print(f"\nDone.")
    print(f"  Extracted:  {n_ok}  ({type_stats})")
    print(f"  Skipped:    {n_skipped}  (doc/ppt need LibreOffice)")
    print(f"  Empty:      {n_empty}")
    print(f"  Failed:     {n_failed}")
    print(f"  FTS index:  {n_indexed} documents searchable")

    db.close()


if __name__ == "__main__":
    main()
